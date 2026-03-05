from __future__ import annotations

import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET


@dataclass
class ExtractionResult:
    text_excerpt: str
    chunk_text: str
    chunk_locations: list[str]
    chunks: list[dict[str, str]]  # [{"location": str, "text": str}]
    content_type: str
    extraction_status: str  # ok | partial | unsupported | error
    metadata: dict[str, Any]


# -------------------------------------------------------------------
# Utilities
# -------------------------------------------------------------------

def _chunks_from_rows(chunk_rows: list[tuple[str, str]]) -> list[dict[str, str]]:
    out: list[dict[str, str]] = []
    for loc, raw in chunk_rows:
        text = (re.sub(r"^\[[^\]]+\]\s*", "", raw).strip() if raw else "") or raw
        out.append({"location": loc, "text": text})
    return out


def _split_chunks(text: str, chunk_size: int = 1200, overlap: int = 150) -> list[str]:
    """Splits text into chunks. overlap=0 avoids duplicating phrases across adjacent chunks."""
    cleaned = re.sub(r"[ \t]+", " ", text or "").strip()
    if not cleaned:
        return []
    if len(cleaned) <= chunk_size:
        return [cleaned]

    step = max(1, chunk_size - overlap)
    out: list[str] = []
    start = 0
    while start < len(cleaned):
        out.append(cleaned[start : start + chunk_size])
        start += step
    return out


def _format_chunks(tag: str, index: int, text: str, overlap: int = 150) -> list[tuple[str, str]]:
    """Split text into chunks. Use overlap=0 for slides/pptx and pdf to avoid same phrase in two chunks."""
    chunks = _split_chunks(text, chunk_size=1200, overlap=overlap)
    if not chunks:
        return []
    if len(chunks) == 1:
        return [(f"{tag}:{index}", f"[{tag} {index}] {chunks[0]}")]
    out: list[tuple[str, str]] = []
    for pos, chunk in enumerate(chunks, start=1):
        out.append((f"{tag}:{index}:{pos}", f"[{tag} {index}.{pos}] {chunk}"))
    return out


def _safe_excerpt(text: str, max_chars: int) -> str:
    if not max_chars:
        return ""
    return (text or "")[:max_chars]


# -------------------------------------------------------------------
# DOCX helpers (pagination)
# -------------------------------------------------------------------

def _read_docx_declared_pages(path: Path) -> int | None:
    """Read document page count from DOCX extended properties (docProps/app.xml)."""
    try:
        with zipfile.ZipFile(path) as archive:
            app_xml = archive.read("docProps/app.xml")
        root = ET.fromstring(app_xml)
        ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
        node = root.find("ep:Pages", ns)
        if node is None or not (node.text or "").strip():
            return None
        pages = int((node.text or "").strip())
        return pages if pages > 0 else None
    except Exception:
        return None


def _choose_docx_break_source(last_rendered_breaks: int, explicit_page_breaks: int) -> tuple[bool, str, int]:
    """Pick DOCX pagination source using real XML markers only.

    Priority:
    1) lastRenderedPageBreak (rendered pagination from Word)
    2) explicit page break (w:br type="page")
    3) estimated fallback
    """
    if last_rendered_breaks > 0:
        return (False, "last_rendered", last_rendered_breaks)
    if explicit_page_breaks > 0:
        return (False, "explicit_break", explicit_page_breaks)
    return (True, "estimated", 0)


def _read_docx_paragraphs_with_breaks(path: Path) -> tuple[list[tuple[str, int, int]], int, int]:
    """Return DOCX paragraphs with per-paragraph break counters from word/document.xml.

    Each item is (paragraph_text, last_rendered_breaks_in_paragraph, explicit_page_breaks_in_paragraph).
    """
    paragraphs: list[tuple[str, int, int]] = []
    last_total = 0
    explicit_total = 0
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(path) as archive:
        document_xml = archive.read("word/document.xml")
    root = ET.fromstring(document_xml)
    body = root.find("w:body", ns)
    if body is None:
        return ([], 0, 0)

    for paragraph in body.findall(".//w:p", ns):
        paragraph_last_rendered = len(paragraph.findall(".//w:lastRenderedPageBreak", ns))
        paragraph_explicit = len(paragraph.findall('.//w:br[@w:type="page"]', ns))
        last_total += paragraph_last_rendered
        explicit_total += paragraph_explicit
        text = "".join((node.text or "") for node in paragraph.findall(".//w:t", ns)).strip()
        paragraphs.append((text, paragraph_last_rendered, paragraph_explicit))

    return (paragraphs, last_total, explicit_total)


# -------------------------------------------------------------------
# Extractors
# -------------------------------------------------------------------

def _extract_plain_text(path: Path, max_chars: int) -> ExtractionResult:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        text = ""
    excerpt = _safe_excerpt(text, max_chars)
    chunks = _format_chunks("section", 1, excerpt)
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text="\n".join(chunk for _, chunk in chunks),
        chunk_locations=[loc for loc, _ in chunks],
        chunks=_chunks_from_rows(chunks),
        content_type="plain_text",
        extraction_status="ok" if excerpt else "partial",
        metadata={"extension": path.suffix.lower()},
    )


def _extract_html(path: Path, max_chars: int) -> ExtractionResult:
    """Best-effort HTML text extraction with zero extra deps."""
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        raw = ""

    title = ""
    m = re.search(r"<title[^>]*>(.*?)</title>", raw, flags=re.IGNORECASE | re.DOTALL)
    if m:
        title = re.sub(r"\s+", " ", m.group(1)).strip()

    cleaned = re.sub(r"<script\b[^>]*>.*?</script>", " ", raw, flags=re.IGNORECASE | re.DOTALL)
    cleaned = re.sub(r"<style\b[^>]*>.*?</style>", " ", cleaned, flags=re.IGNORECASE | re.DOTALL)
    cleaned = re.sub(r"<[^>]+>", " ", cleaned)

    # basic entity cleanup (keep small to avoid needing an HTML parser)
    cleaned = cleaned.replace("&nbsp;", " ").replace("&#160;", " ")
    cleaned = cleaned.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">").replace("&quot;", '"').replace("&apos;", "'")
    cleaned = re.sub(r"\s+", " ", cleaned).strip()

    combined = cleaned
    if title and title.lower() not in combined[:300].lower():
        combined = f"{title}\n{combined}"

    excerpt = _safe_excerpt(combined, max_chars)
    chunks = _format_chunks("section", 1, excerpt)
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text="\n".join(chunk for _, chunk in chunks),
        chunk_locations=[loc for loc, _ in chunks],
        chunks=_chunks_from_rows(chunks),
        content_type="html",
        extraction_status="ok" if excerpt else "partial",
        metadata={"extension": path.suffix.lower(), "title": title},
    )


def _ocr_pdf_page(pdf_path: Path, page_number_1based: int) -> str:
    """Extract text from a PDF page via OCR (pdf2image + Tesseract). Used for scanned pages."""
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        return ""
    try:
        images = convert_from_path(
            str(pdf_path),
            first_page=page_number_1based,
            last_page=page_number_1based,
            dpi=150,
        )
        if not images:
            return ""
        text = pytesseract.image_to_string(images[0], lang="por+eng")
        return (text or "").strip()
    except Exception:
        return ""


def _extract_pdf(path: Path, max_chars: int) -> ExtractionResult:
    from pypdf import PdfReader

    # Optional settings import (keep backwards compatible if app.config isn't present in tests)
    use_ocr = True
    ocr_min_chars = 50
    try:
        from app.config import settings  # type: ignore
        use_ocr = getattr(settings, "pdf_ocr_enabled", True)
        ocr_min_chars = getattr(settings, "pdf_ocr_min_chars", 50)
    except Exception:
        pass

    reader = PdfReader(str(path))
    chunk_rows: list[tuple[str, str]] = []
    pages = len(reader.pages)

    for idx, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        if len(page_text) < ocr_min_chars and use_ocr:
            ocr_text = _ocr_pdf_page(path, idx)
            if ocr_text:
                page_text = ocr_text
        if not page_text:
            continue
        chunk_rows.extend(_format_chunks("page", idx, page_text, overlap=0))
        if max_chars and sum(len(v) for _, v in chunk_rows) >= max_chars * 2:
            break

    chunk_text = "\n".join(value for _, value in chunk_rows)
    excerpt = _safe_excerpt(chunk_text, max_chars)
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text=chunk_text,
        chunk_locations=[loc for loc, _ in chunk_rows],
        chunks=_chunks_from_rows(chunk_rows),
        content_type="pdf",
        extraction_status="ok" if chunk_rows else "partial",
        metadata={"extension": ".pdf", "pages": pages},
    )


def _extract_docx(path: Path, max_chars: int) -> ExtractionResult:
    declared_pages = _read_docx_declared_pages(path)
    chunk_rows: list[tuple[str, str]] = []
    try:
        paragraphs, last_rendered_breaks, explicit_page_breaks = _read_docx_paragraphs_with_breaks(path)
    except Exception:
        paragraphs, last_rendered_breaks, explicit_page_breaks = ([], 0, 0)

    estimated_mode, break_source, selected_breaks = _choose_docx_break_source(
        last_rendered_breaks=last_rendered_breaks,
        explicit_page_breaks=explicit_page_breaks,
    )

    if not paragraphs:
        return ExtractionResult(
            text_excerpt="",
            chunk_text="",
            chunk_locations=[],
            chunks=[],
            content_type="docx",
            extraction_status="partial",
            metadata={
                "extension": ".docx",
                "paragraphs": 0,
                "docx_page_mode": "unavailable",
                "docx_page_markers_found": 0,
                "docx_last_rendered_breaks_found": 0,
                "docx_explicit_page_breaks_found": 0,
                "docx_pages_detected": 0,
                "docx_declared_pages": declared_pages,
            },
        )

    page = 1
    paragraph_in_page = 0
    page_chars = 0
    estimated_chars_per_page = 24000  # fallback when no real pagination markers are present
    paragraphs_with_text = 0

    for paragraph_text, paragraph_last_rendered, paragraph_explicit in paragraphs:
        has_text = bool(paragraph_text.strip())
        if estimated_mode:
            next_len = len(paragraph_text)
            if has_text and paragraph_in_page > 0 and page_chars + next_len > estimated_chars_per_page:
                page += 1
                paragraph_in_page = 0
                page_chars = 0

        if has_text:
            paragraphs_with_text += 1
            paragraph_in_page += 1
            location_prefix = "docx_page_est" if estimated_mode else "docx_page"
            parts = _split_chunks(paragraph_text, chunk_size=1200, overlap=150)
            if len(parts) == 1:
                loc = f"{location_prefix}:{page}:paragraph:{paragraph_in_page}"
                chunk_rows.append((loc, f"[{loc}] {parts[0]}"))
            else:
                for part_idx, part_text in enumerate(parts, start=1):
                    loc = f"{location_prefix}:{page}:paragraph:{paragraph_in_page}:part:{part_idx}"
                    chunk_rows.append((loc, f"[{loc}] {part_text}"))

            if estimated_mode:
                page_chars += len(paragraph_text)

        if not estimated_mode:
            breaks_after = paragraph_last_rendered if break_source == "last_rendered" else paragraph_explicit
            if breaks_after > 0:
                page += breaks_after
                paragraph_in_page = 0

    chunk_text = "\n".join(value for _, value in chunk_rows)
    excerpt = _safe_excerpt(chunk_text, max_chars)

    pages_detected = 0
    if chunk_rows:
        marker = "docx_page_est:" if estimated_mode else "docx_page:"
        pages_detected = max(
            int(loc.split(":")[1])
            for loc, _ in chunk_rows
            if loc.startswith(marker)
        )

    docx_page_mode = break_source if not estimated_mode else "estimated"
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text=chunk_text,
        chunk_locations=[loc for loc, _ in chunk_rows],
        chunks=_chunks_from_rows(chunk_rows),
        content_type="docx",
        extraction_status="ok" if chunk_rows else "partial",
        metadata={
            "extension": ".docx",
            "paragraphs": paragraphs_with_text,
            "docx_page_mode": docx_page_mode,
            "docx_page_markers_found": selected_breaks,
            "docx_last_rendered_breaks_found": last_rendered_breaks,
            "docx_explicit_page_breaks_found": explicit_page_breaks,
            "docx_pages_detected": pages_detected,
            "docx_declared_pages": declared_pages,
        },
    )


def _excel_col_letter(col_idx: int) -> str:
    letters = ""
    n = col_idx
    while n > 0:
        n, rem = divmod(n - 1, 26)
        letters = chr(65 + rem) + letters
    return letters


def _extract_xlsx(path: Path, max_chars: int) -> ExtractionResult:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    chunk_rows: list[tuple[str, str]] = []
    max_cells = 6000
    indexed_cells = 0

    for sheet in wb.worksheets:
        for row_idx, row in enumerate(sheet.iter_rows(values_only=False), start=1):
            for col_idx, cell in enumerate(row, start=1):
                value = getattr(cell, "value", None)
                if value is None:
                    continue
                text = str(value).strip()
                if not text:
                    continue
                col_letters = _excel_col_letter(col_idx)
                marker = f"sheet {sheet.title} row {row_idx} col {col_letters}"
                parts = _split_chunks(text, chunk_size=700, overlap=100)
                if len(parts) == 1:
                    chunk_rows.append((marker, f"[{marker}] {parts[0]}"))
                else:
                    for part_idx, part in enumerate(parts, start=1):
                        loc = f"{marker} part {part_idx}"
                        chunk_rows.append((loc, f"[{loc}] {part}"))
                indexed_cells += 1
                if indexed_cells >= max_cells:
                    break
            if indexed_cells >= max_cells:
                break
        if indexed_cells >= max_cells:
            break

    chunk_text = "\n".join(value for _, value in chunk_rows)
    excerpt = _safe_excerpt(chunk_text, max_chars)
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text=chunk_text,
        chunk_locations=[loc for loc, _ in chunk_rows],
        chunks=_chunks_from_rows(chunk_rows),
        content_type="xlsx",
        extraction_status="ok" if chunk_rows else "partial",
        metadata={"extension": path.suffix.lower(), "sheets": [s.title for s in wb.worksheets]},
    )


def _extract_pptx(path: Path, max_chars: int) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunk_rows: list[tuple[str, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        shapes_sorted = sorted(
            slide.shapes,
            key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)),
        )
        texts: list[str] = []
        for shape in shapes_sorted:
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    row_texts: list[str] = []
                    for cell in row.cells:
                        t = (cell.text_frame.text or "").strip()
                        if t:
                            row_texts.append(t)
                    if row_texts:
                        texts.append(" ".join(row_texts))
            elif hasattr(shape, "text") and shape.text:
                text = str(shape.text).strip()
                if text:
                    texts.append(text)

        if texts:
            chunk_rows.extend(_format_chunks("slide", idx, "\n".join(texts), overlap=0))

    chunk_text = "\n".join(value for _, value in chunk_rows)
    excerpt = _safe_excerpt(chunk_text, max_chars)
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text=chunk_text,
        chunk_locations=[loc for loc, _ in chunk_rows],
        chunks=_chunks_from_rows(chunk_rows),
        content_type="pptx",
        extraction_status="ok" if chunk_rows else "partial",
        metadata={"extension": ".pptx", "slides": len(prs.slides)},
    )


def _extract_legacy_binary(path: Path, max_chars: int) -> ExtractionResult:
    """Safe fallback for legacy Office binaries (.doc/.xls/.ppt): best-effort partial decode."""
    raw = path.read_bytes()
    text = raw.decode("utf-8", errors="ignore")
    if len(text.strip()) < 200:
        text = raw.decode("latin-1", errors="ignore")
    cleaned = "".join(ch if ch.isprintable() else " " for ch in text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    excerpt = _safe_excerpt(cleaned, max_chars)
    chunks = _format_chunks("section", 1, excerpt)
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text="\n".join(chunk for _, chunk in chunks),
        chunk_locations=[loc for loc, _ in chunks],
        chunks=_chunks_from_rows(chunks),
        content_type="legacy_office_binary",
        extraction_status="partial",
        metadata={"extension": path.suffix.lower(), "fallback": "binary_decode"},
    )


def _extract_msg(path: Path, max_chars: int) -> ExtractionResult:
    """Extract subject/body + attachment names from Outlook .msg with extract-msg.

    Notes:
    - We do NOT ingest attachment content by default (avoid ingestion explosion/duplication).
    """
    try:
        import extract_msg  # type: ignore
    except Exception:
        return ExtractionResult(
            text_excerpt="",
            chunk_text="",
            chunk_locations=[],
            chunks=[],
            content_type="msg",
            extraction_status="partial",
            metadata={"extension": ".msg", "fallback": "extract-msg-not-installed"},
        )

    try:
        msg = extract_msg.Message(str(path))
        msg.process()

        subject = (getattr(msg, "subject", "") or "").strip()
        body = (getattr(msg, "body", "") or "").strip()

        attachments: list[str] = []
        for a in (getattr(msg, "attachments", None) or []):
            name = getattr(a, "longFilename", None) or getattr(a, "shortFilename", None) or ""
            name = (name or "").strip()
            if name:
                attachments.append(name)

        combined_parts = []
        if subject:
            combined_parts.append(f"Subject: {subject}")
        if body:
            combined_parts.append(body)
        if attachments:
            combined_parts.append("Attachments: " + ", ".join(attachments))

        combined = "\n".join(combined_parts).strip()
        excerpt = _safe_excerpt(combined, max_chars)
        chunks = _format_chunks("section", 1, excerpt)

        return ExtractionResult(
            text_excerpt=excerpt,
            chunk_text="\n".join(chunk for _, chunk in chunks),
            chunk_locations=[loc for loc, _ in chunks],
            chunks=_chunks_from_rows(chunks),
            content_type="msg",
            extraction_status="ok" if excerpt else "partial",
            metadata={"extension": ".msg", "attachments": attachments[:200]},
        )
    except Exception as exc:
        return ExtractionResult(
            text_excerpt="",
            chunk_text="",
            chunk_locations=[],
            chunks=[],
            content_type="msg",
            extraction_status="error",
            metadata={"extension": ".msg", "error": str(exc)},
        )


def _extract_archive_listing(path: Path, max_chars: int) -> ExtractionResult:
    """Archive ingestion strategy: list contents only (metadata-only), no deep extract by default."""
    ext = path.suffix.lower()
    items: list[str] = []

    if ext == ".zip":
        try:
            with zipfile.ZipFile(path) as z:
                for info in z.infolist()[:5000]:
                    items.append(info.filename)
        except Exception:
            items = []
    elif ext == ".rar":
        # Keep listing empty unless you add 'rarfile' + unrar tool.
        items = []

    text = "\n".join(items)
    excerpt = _safe_excerpt(text, max_chars)
    chunks = _format_chunks("archive_item", 1, excerpt, overlap=0) if excerpt else []
    return ExtractionResult(
        text_excerpt=excerpt,
        chunk_text="\n".join(chunk for _, chunk in chunks),
        chunk_locations=[loc for loc, _ in chunks],
        chunks=_chunks_from_rows(chunks),
        content_type="archive",
        extraction_status="ok" if items else "partial",
        metadata={"extension": ext, "items_listed": len(items)},
    )


# Value used internally when no limit (extraction_mode "all"); extractor processes the entire document.
_EXTRACTION_NO_LIMIT = 100_000_000


def extract_document_content(path: Path, max_chars: int | None = None) -> ExtractionResult:
    """Extract a searchable text representation + chunk locations for citations.

    Supported with good fidelity (your corpus majority):
    - PDF, DOCX, XLSX/XLSM/XLTX, PPTX
    - Plain text containers: .txt/.md/.csv/.json/.log/.eml/.xml/.yaml/.yml
    High-ROI adds (corpus-relevant):
    - HTML/HTM
    - MSG (Outlook)
    Controlled, metadata-only:
    - ZIP/RAR (listing only)
    Fallback:
    - Legacy Office binaries (.doc/.xls/.ppt): best-effort decode, partial
    """
    if max_chars is None:
        max_chars = _EXTRACTION_NO_LIMIT

    ext = path.suffix.lower()
    plain_exts = {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".log",
        ".eml",
        ".xml",
        ".yaml",
        ".yml",
    }

    try:
        if ext in plain_exts:
            return _extract_plain_text(path, max_chars=max_chars)
        if ext in {".html", ".htm"}:
            return _extract_html(path, max_chars=max_chars)
        if ext == ".pdf":
            return _extract_pdf(path, max_chars=max_chars)
        if ext == ".docx":
            return _extract_docx(path, max_chars=max_chars)
        if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
            return _extract_xlsx(path, max_chars=max_chars)
        if ext == ".pptx":
            return _extract_pptx(path, max_chars=max_chars)
        if ext == ".msg":
            return _extract_msg(path, max_chars=max_chars)
        if ext in {".zip", ".rar"}:
            return _extract_archive_listing(path, max_chars=max_chars)
        if ext in {".doc", ".xls", ".ppt"}:
            return _extract_legacy_binary(path, max_chars=max_chars)

        return ExtractionResult(
            text_excerpt="",
            chunk_text="",
            chunk_locations=[],
            chunks=[],
            content_type="unsupported",
            extraction_status="unsupported",
            metadata={"extension": ext},
        )
    except Exception as exc:
        return ExtractionResult(
            text_excerpt="",
            chunk_text="",
            chunk_locations=[],
            chunks=[],
            content_type="error",
            extraction_status="error",
            metadata={"extension": ext, "error": str(exc)},
        )
