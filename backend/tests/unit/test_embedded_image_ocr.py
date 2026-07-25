"""OCR de imagens embutidas em Office (docx/pptx/xlsx) — caso "envelope".

Caso real que motivou a feature (v0.47.0): DOCX com zero runs de texto e o
scan do documento colado como PNG único saía como "sem texto extraível" sem
nunca tentar OCR. Estes testes usam o Tesseract de verdade (mesmo motor do
PDF escaneado e da imagem solta), no padrão de test_image_extraction.py.
"""
from __future__ import annotations

import shutil
import sys
from pathlib import Path

import pytest

from app.document_extractor import extract_document_content

_HAS_TESSERACT = shutil.which("tesseract") is not None


def _make_text_image(path: Path, text: str) -> None:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (900, 200), "white")
    draw = ImageDraw.Draw(img)
    draw.text((20, 60), text, fill="black", font_size=48)
    img.save(path)


def _make_blank_image(path: Path, size: tuple[int, int] = (400, 300)) -> None:
    from PIL import Image

    Image.new("RGB", size, "white").save(path)


def _make_docx_envelope(path: Path, image_paths: list[Path]) -> None:
    """DOCX sem nenhum texto próprio, só imagens inline (o caso real:
    <w:p> com <w:drawing> existe, mas nenhum <w:t> — chunk_rows fica vazio)."""
    from docx import Document

    doc = Document()
    for image_path in image_paths:
        doc.add_picture(str(image_path))
    doc.save(str(path))


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_docx_envelope_extrai_texto_da_imagem_embutida(tmp_path: Path) -> None:
    scan = tmp_path / "scan.png"
    _make_text_image(scan, "ATA DE REUNIAO DO CONSELHO")
    envelope = tmp_path / "envelope.docx"
    _make_docx_envelope(envelope, [scan])

    result = extract_document_content(envelope)
    assert result.content_type == "docx"
    assert result.extraction_status == "ok_ocr"
    assert "ATA" in result.text_excerpt.upper()
    assert result.chunks, "texto OCR deve gerar chunks indexáveis"
    assert result.chunk_locations and result.chunk_locations[0].startswith("image:1")
    assert result.metadata["embedded_images_found"] == 1
    assert result.metadata["embedded_images_ocr"] == 1


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_docx_envelope_com_imagem_ilegivel_segue_partial(tmp_path: Path) -> None:
    blank = tmp_path / "blank.png"
    _make_blank_image(blank)
    envelope = tmp_path / "envelope_ilegivel.docx"
    _make_docx_envelope(envelope, [blank])

    result = extract_document_content(envelope)
    assert result.extraction_status == "partial"
    assert result.text_excerpt == ""
    # Insumo para a UI dizer a causa real ("contém apenas imagem embutida")
    assert result.metadata["embedded_images_found"] == 1
    assert result.metadata["embedded_images_ocr"] == 0


def test_docx_com_texto_proprio_nao_roda_ocr_embutido(tmp_path: Path) -> None:
    """Regra de custo: documento com texto próprio NUNCA paga OCR embutido."""
    from docx import Document

    scan = tmp_path / "scan.png"
    _make_blank_image(scan)
    path = tmp_path / "com_texto.docx"
    doc = Document()
    doc.add_paragraph("Contrato de prestação de serviços entre as partes.")
    doc.add_picture(str(scan))
    doc.save(str(path))

    result = extract_document_content(path)
    assert result.extraction_status == "ok"
    assert "Contrato" in result.text_excerpt
    assert "embedded_images_found" not in result.metadata


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_pptx_envelope_extrai_texto_da_imagem_embutida(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    scan = tmp_path / "diagrama.png"
    _make_text_image(scan, "PLANO DE MIGRACAO")
    path = tmp_path / "deck_envelope.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout blank: sem placeholders de texto
    slide.shapes.add_picture(str(scan), Inches(0), Inches(0))
    prs.save(str(path))

    result = extract_document_content(path)
    assert result.content_type == "pptx"
    assert result.extraction_status == "ok_ocr"
    assert "MIGRACAO" in result.text_excerpt.upper()
    assert result.metadata["embedded_images_found"] == 1


def _make_paragraph_image(path: Path, lines: list[str]) -> None:
    """Imagem com texto longo o suficiente para passar o corte de ruído de 85
    chars (calibrado no corpus real: logo = 0-14 chars, diagrama = 513+)."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (1100, 90 + 70 * len(lines)), "white")
    draw = ImageDraw.Draw(img)
    for i, line in enumerate(lines):
        draw.text((20, 40 + 70 * i), line, fill="black", font_size=44)
    img.save(path)


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_pptx_com_texto_nativo_tambem_le_imagem_de_conteudo(tmp_path: Path) -> None:
    """v0.49.0: no PPTX o OCR roda SEMPRE — deck com texto nativo E diagrama
    colado extrai os dois, com âncora slide:N:image:M."""
    from pptx import Presentation
    from pptx.util import Inches

    diagrama = tmp_path / "diagrama.png"
    _make_paragraph_image(
        diagrama,
        [
            "FLUXO DE INTEGRACAO ENTRE OS SISTEMAS DE COBRANCA",
            "O FATURAMENTO ENVIA OS EVENTOS PARA A FILA CENTRAL",
            "A CONCILIACAO VALIDA OS TOTAIS ANTES DO ENVIO FINAL",
        ],
    )
    path = tmp_path / "deck_misto.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # layout com título
    slide.shapes.title.text = "Arquitetura de Cobrança"
    slide.shapes.add_picture(str(diagrama), Inches(0), Inches(2))
    prs.save(str(path))

    result = extract_document_content(path)
    assert result.extraction_status == "ok"  # texto nativo manda no status
    assert "Arquitetura" in result.text_excerpt
    assert "INTEGRACAO" in result.text_excerpt.upper()
    assert any(loc.startswith("slide:1:image:") for loc in result.chunk_locations)
    assert result.metadata["embedded_images_found"] == 1
    assert result.metadata["embedded_images_ocr"] == 1


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_pptx_logo_com_texto_curto_e_filtrado_como_ruido(tmp_path: Path) -> None:
    """Corte de 85 chars medido no corpus: OCR de logo (texto curto) não vira
    chunk quando o deck tem texto nativo."""
    from pptx import Presentation
    from pptx.util import Inches

    logo = tmp_path / "logo.png"
    _make_text_image(logo, "ACME SA")  # OCR ~7 chars, bem abaixo do corte
    path = tmp_path / "deck_com_logo.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Relatório mensal de operações"
    slide.shapes.add_picture(str(logo), Inches(0), Inches(2))
    prs.save(str(path))

    result = extract_document_content(path)
    assert result.extraction_status == "ok"
    assert not any(":image:" in loc for loc in result.chunk_locations)
    assert result.metadata["embedded_images_found"] == 1
    assert result.metadata["embedded_images_ocr"] == 0


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_xlsx_envelope_extrai_texto_da_imagem_embutida(tmp_path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XlsxImage

    scan = tmp_path / "quadro.png"
    _make_text_image(scan, "ORCAMENTO ANUAL")
    path = tmp_path / "planilha_envelope.xlsx"
    wb = Workbook()
    wb.active.add_image(XlsxImage(str(scan)), "A1")
    wb.save(str(path))

    result = extract_document_content(path)
    assert result.content_type == "xlsx"
    assert result.extraction_status == "ok_ocr"
    assert "ORCAMENTO" in result.text_excerpt.upper()
    assert result.metadata["embedded_images_found"] == 1


@pytest.mark.skipif(not _HAS_TESSERACT, reason="tesseract não instalado no host")
def test_cap_de_imagens_registrado_no_metadata(tmp_path: Path) -> None:
    """12 imagens > cap de 10: nunca silencioso — flag no metadata.

    Dimensões variadas porque o python-docx deduplica imagens byte-idênticas
    em um único media part (fato observado: 12 PNGs iguais → 1 image1.png)."""
    images = []
    for i in range(12):
        img = tmp_path / f"img_{i:02d}.png"
        _make_blank_image(img, size=(60 + i, 40))
        images.append(img)
    envelope = tmp_path / "deck_de_icones.docx"
    _make_docx_envelope(envelope, images)

    result = extract_document_content(envelope)
    assert result.metadata["embedded_images_found"] == 12
    assert result.metadata["embedded_images_ocr_capped"] is True
    assert result.extraction_status == "partial"  # todas em branco: OCR não acha texto


def test_pytesseract_ausente_degrada_sem_erro(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    scan = tmp_path / "scan.png"
    _make_blank_image(scan)
    envelope = tmp_path / "envelope.docx"
    _make_docx_envelope(envelope, [scan])

    # None em sys.modules faz `import pytesseract` levantar ImportError
    monkeypatch.setitem(sys.modules, "pytesseract", None)
    result = extract_document_content(envelope)
    assert result.extraction_status == "partial"
    assert result.metadata["embedded_images_ocr_status"] == "ocr_unavailable"
