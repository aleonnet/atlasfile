"""Unit tests for app.document_extractor (chunks, format, plain text)."""
from __future__ import annotations

import sys
import tempfile
import types
import zipfile
from pathlib import Path

from docx import Document
from docx.enum.text import WD_BREAK
from openpyxl import Workbook
from pptx import Presentation

from app.document_extractor import (
    ExtractionResult,
    _choose_docx_break_source,
    _chunks_from_rows,
    _read_docx_declared_pages,
    _format_chunks,
    _split_chunks,
    extract_document_content,
)


def _write_simple_pdf(path: Path, text: str) -> None:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 18 Tf 50 100 Td ({escaped}) Tj ET".encode("latin-1", errors="ignore")
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (
            b"3 0 obj\n"
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 200] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>\n"
            b"endobj\n"
        ),
        b"4 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        b"5 0 obj\n<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream\nendobj\n",
    ]
    pdf = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for obj in objects:
        offsets.append(len(pdf))
        pdf.extend(obj)
    xref_offset = len(pdf)
    pdf.extend(b"xref\n0 6\n0000000000 65535 f \n")
    for offset in offsets:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n"
        + str(xref_offset).encode("ascii")
        + b"\n%%EOF\n"
    )
    path.write_bytes(bytes(pdf))


def test_split_chunks_empty() -> None:
    assert _split_chunks("") == []
    assert _split_chunks("   ") == []


def test_split_chunks_short() -> None:
    text = "one two three"
    assert _split_chunks(text, chunk_size=1200) == [text]


def test_split_chunks_long() -> None:
    text = "x " * 800  # > 1200 chars
    chunks = _split_chunks(text, chunk_size=1200, overlap=150)
    assert len(chunks) >= 2
    for c in chunks:
        assert len(c) <= 1200 + 50


def test_format_chunks_single() -> None:
    out = _format_chunks("page", 1, "short text", overlap=0)
    assert len(out) == 1
    assert out[0][0] == "page:1"
    assert "[page 1]" in out[0][1]


def test_format_chunks_multiple() -> None:
    long_text = "word " * 400  # multiple chunks
    out = _format_chunks("page", 2, long_text, overlap=150)
    assert len(out) >= 2
    assert out[0][0].startswith("page:2:1")
    assert "[page 2.1]" in out[0][1]


def test_chunks_from_rows() -> None:
    rows = [("page:1", "[page 1] hello"), ("page:2", "[page 2] world")]
    out = _chunks_from_rows(rows)
    assert len(out) == 2
    assert out[0]["location"] == "page:1"
    assert "hello" in out[0]["text"]
    assert out[1]["location"] == "page:2"
    assert "world" in out[1]["text"]


def test_extract_plain_text_txt_file() -> None:
    with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
        f.write("Hello world\nLine two")
        path = Path(f.name)
    try:
        result = extract_document_content(path, max_chars=1000)
        assert isinstance(result, ExtractionResult)
        assert "Hello" in result.text_excerpt
        assert result.extraction_status == "ok"
        assert result.content_type == "plain_text"
    finally:
        path.unlink(missing_ok=True)


def test_extract_html_file() -> None:
    html = """
    <html>
      <head>
        <title>Ata de Reuniao</title>
        <style>.x { color: red; }</style>
      </head>
      <body>
        <script>console.log("ignore")</script>
        <h1>Conteudo principal</h1>
      </body>
    </html>
    """
    with tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8") as f:
        f.write(html)
        path = Path(f.name)
    try:
        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "html"
        assert result.extraction_status == "ok"
        assert "Ata de Reuniao" in result.text_excerpt
        assert "Conteudo principal" in result.text_excerpt
        assert "console.log" not in result.text_excerpt
    finally:
        path.unlink(missing_ok=True)


def test_extract_archive_zip_listing() -> None:
    with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as f:
        path = Path(f.name)
    try:
        with zipfile.ZipFile(path, "w") as archive:
            archive.writestr("pasta/doc1.txt", "abc")
            archive.writestr("doc2.pdf", "xyz")
        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "archive"
        assert result.extraction_status == "ok"
        assert "pasta/doc1.txt" in result.text_excerpt
        assert "doc2.pdf" in result.text_excerpt
        assert result.metadata["items_listed"] == 2
    finally:
        path.unlink(missing_ok=True)


def test_extract_archive_rar_listing_with_stubbed_parser(monkeypatch) -> None:
    fake_module = types.ModuleType("rarfile")

    class _Info:
        def __init__(self, filename: str) -> None:
            self.filename = filename

    class _RarFile:
        def __init__(self, _path: Path) -> None:
            self._items = [_Info("pasta/doc1.txt"), _Info("doc2.pdf")]

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return False

        def infolist(self):
            return self._items

    fake_module.RarFile = _RarFile
    monkeypatch.setitem(sys.modules, "rarfile", fake_module)

    with tempfile.NamedTemporaryFile(suffix=".rar", delete=False) as f:
        path = Path(f.name)
    try:
        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "archive"
        assert result.extraction_status == "ok"
        assert "pasta/doc1.txt" in result.text_excerpt
        assert "doc2.pdf" in result.text_excerpt
        assert result.metadata["items_listed"] == 2
    finally:
        path.unlink(missing_ok=True)


def test_extract_plain_text_xml_extension() -> None:
    with tempfile.NamedTemporaryFile(mode="w", suffix=".xml", delete=False, encoding="utf-8") as f:
        f.write("<root><name>Atlas</name></root>")
        path = Path(f.name)
    try:
        result = extract_document_content(path, max_chars=200)
        assert result.content_type == "plain_text"
        assert result.extraction_status == "ok"
        assert "<name>Atlas</name>" in result.text_excerpt
    finally:
        path.unlink(missing_ok=True)


def test_extract_msg_with_stubbed_module(monkeypatch) -> None:
    fake_module = types.ModuleType("extract_msg")

    class _Attachment:
        def __init__(self, name: str) -> None:
            self.longFilename = name
            self.shortFilename = ""

    class _Message:
        def __init__(self, _path: str) -> None:
            self.subject = "Assunto Teste"
            self.body = "Corpo Teste"
            self.attachments = [_Attachment("arquivo.pdf")]

        def process(self) -> None:
            return None

    fake_module.Message = _Message
    monkeypatch.setitem(sys.modules, "extract_msg", fake_module)

    with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as f:
        path = Path(f.name)
    try:
        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "msg"
        assert result.extraction_status == "ok"
        assert "Assunto Teste" in result.text_excerpt
        assert "Corpo Teste" in result.text_excerpt
        assert "arquivo.pdf" in result.text_excerpt
    finally:
        path.unlink(missing_ok=True)


def test_extract_pdf_from_generated_real_file() -> None:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        path = Path(f.name)
    try:
        _write_simple_pdf(path, "Atlas PDF Contrato")
        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "pdf"
        assert result.extraction_status == "ok"
        assert "Atlas PDF Contrato" in result.text_excerpt
        assert "page:1" in result.chunk_locations
    finally:
        path.unlink(missing_ok=True)


def test_extract_pptx_from_generated_real_file() -> None:
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = Path(f.name)
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(left=0, top=0, width=7_000_000, height=500_000)
        title_box.text_frame.text = "Kickoff Neptune"
        table_shape = slide.shapes.add_table(rows=2, cols=2, left=0, top=700_000, width=7_000_000, height=1_200_000)
        table = table_shape.table
        table.cell(0, 0).text = "Tema"
        table.cell(0, 1).text = "Status"
        table.cell(1, 0).text = "Cronograma"
        table.cell(1, 1).text = "Aprovado"
        prs.save(str(path))

        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "pptx"
        assert result.extraction_status == "ok"
        assert "Kickoff Neptune" in result.text_excerpt
        assert "Cronograma Aprovado" in result.text_excerpt
        assert "slide:1" in result.chunk_locations
    finally:
        path.unlink(missing_ok=True)


def test_extract_xlsx_from_generated_real_file() -> None:
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        path = Path(f.name)
    try:
        wb = Workbook()
        ws = wb.active
        ws.title = "Financeiro"
        ws["A1"] = "Conta"
        ws["B1"] = "Valor"
        ws["A2"] = "Caixa"
        ws["B2"] = 1250
        wb.save(str(path))

        result = extract_document_content(path, max_chars=1000)
        assert result.content_type == "xlsx"
        assert result.extraction_status == "ok"
        assert "sheet Financeiro row 1 col A" in result.chunk_text
        assert "Caixa" in result.text_excerpt
        assert "1250" in result.text_excerpt
    finally:
        path.unlink(missing_ok=True)


def test_extract_msg_from_real_fixture() -> None:
    fixture = Path(__file__).resolve().parents[1] / "fixtures" / "msg" / "strangeDate.msg"
    result = extract_document_content(fixture, max_chars=1500)

    assert result.content_type == "msg"
    assert result.extraction_status == "ok"
    assert "MSG Test File" in result.text_excerpt
    assert "Provide example of this file type" in result.text_excerpt


def test_extract_docx_locations_with_explicit_page_break() -> None:
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        path = Path(f.name)
    try:
        doc = Document()
        p1 = doc.add_paragraph("Primeiro paragrafo do documento de teste.")
        p1.add_run().add_break(WD_BREAK.PAGE)
        doc.add_paragraph("Segundo paragrafo apos quebra de pagina.")
        doc.save(str(path))

        result = extract_document_content(path, max_chars=100000)
        assert result.content_type == "docx"
        assert result.extraction_status == "ok"
        assert "docx_page:1:paragraph:1" in result.chunk_locations
        assert "docx_page:2:paragraph:1" in result.chunk_locations
        assert result.metadata["docx_page_mode"] == "explicit_break"
    finally:
        path.unlink(missing_ok=True)


def test_extract_docx_locations_without_page_markers_use_estimated_mode() -> None:
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        path = Path(f.name)
    try:
        doc = Document()
        doc.add_paragraph("Paragrafo um sem quebra explicita.")
        doc.add_paragraph("Paragrafo dois sem quebra explicita.")
        doc.save(str(path))

        result = extract_document_content(path, max_chars=100000)
        assert result.content_type == "docx"
        assert result.chunk_locations
        assert all(loc.startswith("docx_page_est:") for loc in result.chunk_locations)
        assert result.metadata["docx_page_mode"] == "estimated"
    finally:
        path.unlink(missing_ok=True)


def test_choose_docx_break_source_prioritizes_last_rendered() -> None:
    estimated, source, selected = _choose_docx_break_source(last_rendered_breaks=194, explicit_page_breaks=10)
    assert estimated is False
    assert source == "last_rendered"
    assert selected == 194


def test_choose_docx_break_source_uses_explicit_when_last_missing() -> None:
    estimated, source, selected = _choose_docx_break_source(last_rendered_breaks=0, explicit_page_breaks=3)
    assert estimated is False
    assert source == "explicit_break"
    assert selected == 3


def test_choose_docx_break_source_falls_back_to_estimated() -> None:
    estimated, source, selected = _choose_docx_break_source(last_rendered_breaks=0, explicit_page_breaks=0)
    assert estimated is True
    assert source == "estimated"
    assert selected == 0


def test_read_docx_declared_pages_from_generated_file() -> None:
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        path = Path(f.name)
    try:
        doc = Document()
        doc.add_paragraph("Teste.")
        doc.save(str(path))
        pages = _read_docx_declared_pages(path)
        assert pages is None or pages > 0
    finally:
        path.unlink(missing_ok=True)
